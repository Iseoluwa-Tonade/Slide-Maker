import streamlit as st
import os
import json
import time
import requests
import io
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.text import PP_ALIGN

# --- Configuration ---
OUTPUT_FILENAME = "edited_presentation.pptx"

# --- LLM API Configuration ---
API_KEY = os.environ.get("GEMINI_API_KEY")
API_URL_JSON = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-preview-09-2025:generateContent?key={API_KEY}"

# --- CORE PPTX LOGIC ---

def apply_text_formatting(target_run, source_font):
    """Reapplies font size, name, bold, and color from source_font to target_run."""
    if source_font:
        if source_font.size: target_run.font.size = source_font.size
        if source_font.name: target_run.font.name = source_font.name
        target_run.font.bold = source_font.bold

        color_type = source_font.color.type
        if color_type == MSO_COLOR_TYPE.RGB:
            target_run.font.color.rgb = source_font.color.rgb
        elif color_type == MSO_COLOR_TYPE.SCHEME:
            target_run.font.color.theme_color = source_font.color.theme_color
            if hasattr(source_font.color, 'brightness') and source_font.color.brightness is not None:
                target_run.font.color.brightness = source_font.color.brightness

def replace_text_safely(shape, new_text):
    """Replaces text in a shape while preserving the formatting of the first run."""
    text_frame = shape.text_frame

    if not text_frame.paragraphs:
        p0 = text_frame.add_paragraph()
    else:
        p0 = text_frame.paragraphs[0]

    r0 = p0.runs[0].font if p0.runs else None
    original_alignment = p0.alignment
    original_level = p0.level

    while len(text_frame.paragraphs) > 1:
        p = text_frame.paragraphs[-1]
        if hasattr(p, '_element'): p._element.getparent().remove(p._element)
        else: break
    p0.clear()

    new_lines = new_text.split('\n')

    if new_lines:
        run = p0.add_run()
        run.text = new_lines[0]
        p0.alignment = original_alignment
        p0.level = original_level
        apply_text_formatting(run, r0)

        for line in new_lines[1:]:
            p = text_frame.add_paragraph()
            p.text = line
            p.level = original_level
            p.alignment = original_alignment
            if p.runs: apply_text_formatting(p.runs[0], r0)

def extract_presentation_text_blocks(prs):
    """Scans the presentation and returns a list of text blocks and a shape map."""
    presentation_blocks = []
    shape_map = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                current_text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if current_text.strip():
                    presentation_blocks.append({
                        "id": len(presentation_blocks),
                        "slide_number": slide_idx + 1,
                        "current_text": current_text,
                        "shape_name": shape.name,
                    })
                    shape_map.append(shape)
    return presentation_blocks, shape_map

# --- CORE AI LOGIC ---

def call_gemini_api(url, payload):
    """Helper function to call the Gemini API with retry logic."""
    max_retries = 3
    delay = 2 
    for attempt in range(max_retries):
        try:
            response = requests.post(url, headers={'Content-Type': 'application/json'}, data=json.dumps(payload))
            response.raise_for_status()
            result = response.json()
            if 'candidates' in result and result['candidates']:
                return result['candidates'][0]['content']['parts'][0]['text']
            return None
        except requests.exceptions.HTTPError as e:
            st.error(f"API Error (Attempt {attempt + 1}): {e.response.status_code}")
            if attempt < max_retries - 1:
                time.sleep(delay)
                delay *= 2
            else:
                st.error("Max retries reached. Aborting AI call.")
                raise
        except Exception as e:
            st.error(f"An unexpected error occurred during API call: {e}")
            raise
    return None

@st.cache_data(show_spinner=False)
def generate_ai_mapping(_transcript, presentation_blocks, api_key):
    """Calls Gemini to map transcript data to text blocks (used for Batch Mode)."""
    if not api_key: return []
    url = f"https{API_URL_JSON.split('https')[-1]}" # Rebuild URL to ensure correctness

    MAX_TRANSCRIPT_LENGTH = 8000
    transcript_to_send = _transcript
    if len(_transcript) > MAX_TRANSCRIPT_LENGTH:
        head = _transcript[:4000]
        tail = _transcript[-4000:]
        transcript_to_send = f"{head}\n\n[... TRUNCATED ...]\n\n{tail}"

    schema = {
        "type": "ARRAY",
        "description": "A list of text replacements. Only include blocks that require an update based on the transcript.",
        "items": {
            "type": "OBJECT",
            "properties": {
                "original_index": {"type": "NUMBER", "description": "The 'id' number of the text block to replace from the PRESENTATION TEXT BLOCKS list."},
                "new_text": {"type": "STRING", "description": "The concise, new text or bulleted content derived from the transcript. Use '\\n' for new lines/bullets."}
            },
            "required": ["original_index", "new_text"]
        }
    }

    system_prompt = "You are a Presentation Automation Specialist. Analyze the TRANSCRIPT and the PRESENTATION TEXT BLOCKS. Only create entries for blocks that need modification."
    presentation_list_text = json.dumps(presentation_blocks, indent=2)
    user_query = f"MEETING TRANSCRIPT:\n---\n{transcript_to_send}\n---\n\NPRESENTATION TEXT BLOCKS (Index to Map to):\n---\n{presentation_list_text}\n---\nGenerate a JSON array of objects mapping the 'original_index' to the derived 'new_text'."

    payload = {
        "contents": [{"parts": [{"text": user_query}]}],
        "systemInstruction": {"parts": [{"text": system_prompt}]},
        "generationConfig": {"responseMimeType": "application/json", "responseSchema": schema}
    }

    with st.spinner("Analyzing transcript and mapping changes with Gemini..."):
        json_text = call_gemini_api(url, payload)
    if json_text: return json.loads(json_text)
    return []

def process_batch_edits(prs, edited_block_lines, shape_map):
    """Parses the bulk edited text block and applies changes."""
    applied_count = 0
    for line in edited_block_lines:
        try:
            match = line.split('|', 3)
            if len(match) < 4: continue

            id_str = match[0].strip()[1:]
            block_id = int(id_str.split('|')[0])
            new_text_raw = match[3].strip()
            new_text = new_text_raw.replace(' \\n ', '\n').replace('\\n', '\n')

            if 0 <= block_id < len(shape_map):
                shape = shape_map[block_id]
                replace_text_safely(shape, new_text)
                applied_count += 1
        except Exception as e:
            st.warning(f"Could not parse or apply change for line: {line}. Skipping. Error: {e}")
    return applied_count

# --- STREAMLIT UI ---

def main_app():
    st.set_page_config(layout="wide", page_title="PPTX AI Co-Editor")

    # --- Initialize Session State ---
    if 'editing_started' not in st.session_state: st.session_state['editing_started'] = False
    if 'editor_text' not in st.session_state: st.session_state['editor_text'] = ""

    st.title("📄 AI PowerPoint (PPTX) Automation Tool")
    st.markdown("Upload your **.pptx** template and paste a meeting transcript to automatically update the text.")

    if not API_KEY:
        st.error("⚠️ **Deployment Error:** The `GEMINI_API_KEY` environment variable is not set.")

    # 1. Inputs
    col1, col2 = st.columns([1, 1])
    with col1:
        uploaded_file = st.file_uploader("Upload PowerPoint File (.pptx)", type="pptx", key='uploaded_file')
    with col2:
        transcript = st.text_area("Paste Meeting Transcript Here", height=300, key='transcript_area')

    st.markdown("---")

    # 2. Execution Button
    if st.button("Start Analysis and Editing 🚀", disabled=(not uploaded_file or not transcript or not API_KEY)):
        try:
            prs = Presentation(io.BytesIO(uploaded_file.getvalue()))
        except Exception as e:
            st.error(f"Could not load presentation: {e}")
            return

        presentation_blocks, shape_map = extract_presentation_text_blocks(prs)
        if not presentation_blocks:
            st.warning("No editable text blocks found in this PPTX file.")
            return

        ai_mapping_list = generate_ai_mapping(transcript, presentation_blocks, API_KEY)
        ai_replacements = {item['original_index']: item['new_text'] for item in ai_mapping_list}

        # --- Generate the Comprehensive Edit Block ---
        edit_block_output = []
        for block in presentation_blocks:
            original_index = block['id']
            slide_num = block['slide_number']
            indicator = "AI_SUGGESTED" if original_index in ai_replacements else "ORIGINAL_TEXT"
            initial_text = ai_replacements.get(original_index, block['current_text']).replace("\n", " \\n ")
            formatted_line = f"[{original_index}|S{slide_num}|{indicator}] | {initial_text}"
            edit_block_output.append(formatted_line)
        
        st.session_state['editor_text'] = "\n".join(edit_block_output)
        st.session_state['prs_buffer'] = io.BytesIO(uploaded_file.getvalue()) # Store original buffer
        st.session_state['shape_map'] = shape_map
        st.session_state['editing_started'] = True
        
        st.success("Analysis complete. Scroll down to the Batch Editor.")
        st.rerun()

    # --- Editing Interface (Batch Mode) ---
    if st.session_state['editing_started']:
        st.subheader("📝 Batch Editor: All Slide Text")
        st.markdown(
            """
            Review and edit all text below. The AI suggestions have been pre-filled.
            - **Edit ONLY** the text after the last `|` delimiter.
            - Use `\\n` for new lines or bullet points.
            """
        )

        edited_block_text = st.text_area(
            "Copy, Edit, and Paste Back (if needed):",
            value=st.session_state['editor_text'],
            height=600,
            key='batch_edit_area'
        )
        
        # Update session state as user types
        st.session_state['editor_text'] = edited_block_text

        st.markdown("---")
        
        col_preview, col_final = st.columns(2)

        with col_preview:
            if st.button("Generate Preview File", help="Apply changes and download a copy to check your work."):
                prs = Presentation(st.session_state['prs_buffer']) # Re-load original
                shape_map = st.session_state['shape_map']
                edited_block_lines = st.session_state['editor_text'].split('\n')
                
                with st.spinner("Applying changes for preview..."):
                    process_batch_edits(prs, edited_block_lines, shape_map)

                output_buffer = io.BytesIO()
                prs.save(output_buffer)
                output_buffer.seek(0)

                st.download_button(
                    label="⬇️ Download Preview File",
                    data=output_buffer,
                    file_name="preview_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key='download_preview'
                )

        with col_final:
            if st.button("✅ Apply Final Changes and Download", type="primary"):
                prs = Presentation(st.session_state['prs_buffer']) # Re-load original
                shape_map = st.session_state['shape_map']
                edited_block_lines = st.session_state['editor_text'].split('\n')
                
                with st.spinner("Applying final changes..."):
                    process_batch_edits(prs, edited_block_lines, shape_map)

                output_buffer = io.BytesIO()
                prs.save(output_buffer)
                output_buffer.seek(0)
                
                st.success(f"✅ Success! Your file is ready.")
                st.download_button(
                    label=f"⬇️ Download Final {OUTPUT_FILENAME}",
                    data=output_buffer,
                    file_name=OUTPUT_FILENAME,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key='download_final'
                )

if __name__ == "__main__":
    main_app()
